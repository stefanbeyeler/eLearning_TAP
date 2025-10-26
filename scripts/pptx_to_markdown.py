#!/usr/bin/env python3
"""Convert PowerPoint presentation to Markdown"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import sys
import os

def extract_text_from_shape(shape):
    """Extract text from a shape, handling different types"""
    text_parts = []

    if hasattr(shape, "text") and shape.text:
        text_parts.append(shape.text)

    if hasattr(shape, "table"):
        # Handle tables
        table = shape.table
        table_text = []
        for row in table.rows:
            row_text = " | ".join([cell.text for cell in row.cells])
            table_text.append(row_text)
        if table_text:
            text_parts.append("\n".join(table_text))

    return "\n".join(text_parts)

def extract_images_from_slide(slide, slide_num, images_dir, images_md_path=None, use_html_syntax=True):
    """Extract images from a slide and return markdown references"""
    image_refs = []
    image_counter = 1

    # Use images_md_path for markdown references, or fall back to images_dir
    md_path = images_md_path if images_md_path is not None else images_dir

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob

                # Create filename with .webp extension
                filename = f"slide{slide_num:02d}_img{image_counter}.webp"
                filepath = os.path.join(images_dir, filename)

                # Convert to WebP using PIL
                img = Image.open(io.BytesIO(image_bytes))
                # Convert to RGB if necessary (WebP doesn't support RGBA with high quality)
                if img.mode in ('RGBA', 'LA', 'P'):
                    # Create white background
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')

                # Save as WebP with quality 85
                img.save(filepath, 'WEBP', quality=85)

                # Create markdown reference with HTML build syntax or standard markdown
                if use_html_syntax:
                    image_refs.append(f"**BILD: {md_path}/{filename}**")
                else:
                    image_refs.append(f"![Bild {image_counter}]({md_path}/{filename})")
                image_counter += 1

            except Exception as e:
                print(f"Warning: Could not extract image from slide {slide_num}: {e}")

        # Check for images in group shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = sub_shape.image
                        image_bytes = image.blob

                        # Create filename with .webp extension
                        filename = f"slide{slide_num:02d}_img{image_counter}.webp"
                        filepath = os.path.join(images_dir, filename)

                        # Convert to WebP using PIL
                        img = Image.open(io.BytesIO(image_bytes))
                        if img.mode in ('RGBA', 'LA', 'P'):
                            background = Image.new('RGB', img.size, (255, 255, 255))
                            if img.mode == 'P':
                                img = img.convert('RGBA')
                            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                            img = background
                        elif img.mode != 'RGB':
                            img = img.convert('RGB')

                        img.save(filepath, 'WEBP', quality=85)

                        if use_html_syntax:
                            image_refs.append(f"**BILD: {md_path}/{filename}**")
                        else:
                            image_refs.append(f"![Bild {image_counter}]({md_path}/{filename})")
                        image_counter += 1

                    except Exception as e:
                        print(f"Warning: Could not extract grouped image from slide {slide_num}: {e}")

    return image_refs

def pptx_to_markdown(pptx_path, images_dir="pictures", images_md_path=None, use_html_syntax=True):
    """Convert PowerPoint presentation to Markdown format

    Args:
        pptx_path: Path to PowerPoint file
        images_dir: Physical directory to save images
        images_md_path: Path to use in markdown references (defaults to images_dir)
        use_html_syntax: Use **BILD:** syntax if True, else use ![alt](path)
    """
    prs = Presentation(pptx_path)
    markdown_content = []

    # Create images directory if it doesn't exist
    if not os.path.exists(images_dir):
        os.makedirs(images_dir)
        print(f"Created directory: {images_dir}")

    # Add title
    markdown_content.append("# Schulung Grundlagenmodul\n")

    for slide_num, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        if title:
            markdown_content.append(f"\n## Folie {slide_num}: {title}\n")
        else:
            markdown_content.append(f"\n## Folie {slide_num}\n")

        # Extract images first
        image_refs = extract_images_from_slide(slide, slide_num, images_dir, images_md_path, use_html_syntax)
        if image_refs:
            for img_ref in image_refs:
                markdown_content.append(img_ref)
            markdown_content.append("")  # Add spacing after images

        # Extract text from all shapes
        for shape in slide.shapes:
            # Skip the title shape as we already processed it
            if shape == slide.shapes.title:
                continue

            # Skip picture shapes (already handled)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue

            text = extract_text_from_shape(shape)
            if text:
                # Check if it's a bullet point or regular text
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        # Add bullet points for indented or list items
                        if line.startswith(('\t', '  ', '•', '-', '·')):
                            cleaned_line = line.strip('•-·\t ')
                            markdown_content.append(f"- {cleaned_line}")
                        else:
                            markdown_content.append(f"{line}")

                markdown_content.append("")  # Add spacing

        # Add separator between slides
        markdown_content.append("---\n")

    return "\n".join(markdown_content)

if __name__ == "__main__":
    pptx_path = "documents/Schulung Grundlagenmodul[52].pptx"
    output_path = "Schulung_Grundlagenmodul.md"

    try:
        # Convert with HTML syntax for build_html.py compatibility
        markdown_text = pptx_to_markdown(pptx_path, images_dir="content/images", use_html_syntax=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)

        print(f"Successfully converted {pptx_path} to {output_path}")
        print(f"Total length: {len(markdown_text)} characters")
        print(f"\nImages saved to: content/images/")
        print(f"Using **BILD: path** syntax for build_html.py compatibility")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
